from __future__ import annotations

import json
from pathlib import Path

from lxml import etree
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_DIR = ROOT / "test-data" / "pptx-fixtures"
ASSET_DIR = ROOT / "test-data" / "assets"
MANIFEST_PATH = ROOT / "test-data" / "fixtures-manifest.json"


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def create_sample_image() -> Path:
    image_path = ASSET_DIR / "sample-image.png"
    img = Image.new("RGB", (240, 140), color=(240, 240, 240))
    draw = ImageDraw.Draw(img)

    # Simple deterministic pattern for image parser validation.
    draw.rectangle((10, 10, 230, 130), outline=(31, 78, 121), width=4)
    draw.ellipse((40, 30, 110, 100), fill=(212, 80, 61), outline=(100, 20, 20), width=3)
    draw.rectangle((130, 35, 210, 105), fill=(76, 163, 74), outline=(30, 80, 30), width=3)

    img.save(image_path)
    return image_path


def set_slide_background(slide, rgb: tuple[int, int, int]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)


def add_notes(slide, text: str) -> None:
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    notes_tf.text = text


def make_fixture_01_basic_text_shapes() -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (255, 255, 255))

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8.6), Inches(0.7))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Fixture 01: Basic Text and Shapes"
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = RGBColor(208, 68, 35)

    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(8.2), Inches(1.2))
    tf = body_box.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "Bold + italic + aligned text"
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.runs[0]
    r1.font.bold = True
    r1.font.italic = True
    r1.font.size = Pt(18)
    r1.font.color.rgb = RGBColor(50, 50, 50)

    p2 = tf.add_paragraph()
    p2.text = "Center aligned second paragraph"
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.runs[0]
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(20, 90, 155)

    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.1), Inches(2.8), Inches(1.2))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(172, 212, 51)
    rect.line.color.rgb = RGBColor(80, 120, 20)
    rect.line.width = Pt(2)

    ell = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.1), Inches(3.0), Inches(2.0), Inches(1.4))
    ell.fill.solid()
    ell.fill.fore_color.rgb = RGBColor(92, 184, 212)
    ell.line.color.rgb = RGBColor(31, 78, 121)
    ell.line.width = Pt(2)

    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(6.7), Inches(3.2), Inches(8.8), Inches(4.1))
    line.line.color.rgb = RGBColor(208, 68, 35)
    line.line.width = Pt(3)

    add_notes(slide, "Fixture 01 note text")

    out_name = "fixture-01-basic-text-shapes.pptx"
    prs.save(OUT_DIR / out_name)
    return out_name


def make_fixture_02_image_table(image_path: Path) -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (247, 247, 247))

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(8.8), Inches(0.6))
    p = title_box.text_frame.paragraphs[0]
    p.text = "Fixture 02: Image + Table"
    p.alignment = PP_ALIGN.LEFT
    r = p.runs[0]
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = RGBColor(31, 78, 121)

    slide.shapes.add_picture(str(image_path), Inches(0.8), Inches(1.2), width=Inches(4.0), height=Inches(2.3))

    rows, cols = 3, 3
    table_shape = slide.shapes.add_table(rows, cols, Inches(5.1), Inches(1.2), Inches(4.0), Inches(2.3))
    table = table_shape.table

    headers = ["Name", "Type", "Value"]
    values = [
        ["A", "alpha", "10"],
        ["B", "beta", "20"],
    ]

    for c in range(cols):
        table.cell(0, c).text = headers[c]
        table.cell(0, c).fill.solid()
        table.cell(0, c).fill.fore_color.rgb = RGBColor(220, 235, 247)

    for r_i, row_vals in enumerate(values, start=1):
        for c_i, val in enumerate(row_vals):
            table.cell(r_i, c_i).text = val

    add_notes(slide, "Fixture 02 note text")

    out_name = "fixture-02-image-table.pptx"
    prs.save(OUT_DIR / out_name)
    return out_name


def make_fixture_03_multislide_notes() -> str:
    prs = Presentation()

    # Slide 1
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s1, (255, 255, 255))
    b1 = s1.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.3), Inches(1.0))
    p1 = b1.text_frame.paragraphs[0]
    p1.text = "Slide 1 - Left"
    p1.alignment = PP_ALIGN.LEFT
    p1.runs[0].font.size = Pt(28)
    add_notes(s1, "Fixture 03 - notes for slide 1")

    # Slide 2
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s2, (235, 245, 255))
    b2 = s2.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.3), Inches(1.0))
    p2 = b2.text_frame.paragraphs[0]
    p2.text = "Slide 2 - Center"
    p2.alignment = PP_ALIGN.CENTER
    p2.runs[0].font.size = Pt(30)
    p2.runs[0].font.color.rgb = RGBColor(31, 78, 121)
    add_notes(s2, "Fixture 03 - notes for slide 2")

    # Slide 3
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s3, (245, 235, 235))
    b3 = s3.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(8.3), Inches(1.0))
    p3 = b3.text_frame.paragraphs[0]
    p3.text = "Slide 3 - Right"
    p3.alignment = PP_ALIGN.RIGHT
    p3.runs[0].font.size = Pt(30)
    p3.runs[0].font.bold = True
    p3.runs[0].font.color.rgb = RGBColor(140, 40, 40)
    add_notes(s3, "Fixture 03 - notes for slide 3")

    out_name = "fixture-03-multi-slide-notes.pptx"
    prs.save(OUT_DIR / out_name)
    return out_name


def make_fixture_04_bullets_rotation() -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (255, 255, 255))

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.7), Inches(0.7))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Fixture 04: Bullets and Rotation"
    tp.runs[0].font.size = Pt(24)
    tp.runs[0].font.bold = True

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(4.4), Inches(2.5))
    tf = body.text_frame
    tf.text = "Level 0 bullet"
    tf.paragraphs[0].level = 0
    p = tf.add_paragraph()
    p.text = "Level 1 bullet"
    p.level = 1
    p2 = tf.add_paragraph()
    p2.text = "Level 2 bullet"
    p2.level = 2

    rot = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(5.8), Inches(1.8), Inches(2.8), Inches(1.4))
    rot.rotation = 25
    rot.fill.solid()
    rot.fill.fore_color.rgb = RGBColor(255, 191, 0)
    rot.line.color.rgb = RGBColor(130, 90, 0)

    add_notes(slide, "Fixture 04 note text")

    out_name = "fixture-04-bullets-rotation.pptx"
    prs.save(OUT_DIR / out_name)
    return out_name


def make_fixture_05_group_transform() -> str:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, (255, 255, 255))

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.7), Inches(0.7))
    tp = title.text_frame.paragraphs[0]
    tp.text = "Fixture 05: Group Transform"
    tp.runs[0].font.size = Pt(24)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = RGBColor(66, 66, 66)

    group = slide.shapes.add_group_shape()

    # Child shapes are defined in group-local coordinates.
    r = group.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(0.2), Inches(1.8), Inches(1.0))
    r.fill.solid()
    r.fill.fore_color.rgb = RGBColor(106, 168, 79)
    r.line.color.rgb = RGBColor(56, 118, 29)
    r.line.width = Pt(1.5)

    e = group.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.3), Inches(0.25), Inches(1.2), Inches(0.9))
    e.fill.solid()
    e.fill.fore_color.rgb = RGBColor(61, 133, 198)
    e.line.color.rgb = RGBColor(31, 78, 121)
    e.line.width = Pt(1.5)

    c = group.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.4), Inches(1.5), Inches(3.3), Inches(1.85))
    c.line.color.rgb = RGBColor(204, 102, 0)
    c.line.width = Pt(2)

    # Nested group for recursive grpSp transform checks.
    nested = group.shapes.add_group_shape()
    n1 = nested.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(0.0), Inches(0.0), Inches(1.0), Inches(0.6))
    n1.fill.solid()
    n1.fill.fore_color.rgb = RGBColor(255, 192, 0)
    n1.line.color.rgb = RGBColor(153, 115, 0)
    n1.rotation = 15
    n2 = nested.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.05), Inches(0.05), Inches(0.55), Inches(0.45))
    n2.fill.solid()
    n2.fill.fore_color.rgb = RGBColor(142, 124, 195)
    n2.line.color.rgb = RGBColor(103, 78, 167)

    # Position and scale nested group in parent group local space.
    nested.left = Inches(1.35)
    nested.top = Inches(1.95)
    nested.width = Inches(2.1)
    nested.height = Inches(0.9)

    # Position and scale outer group on slide.
    group.left = Inches(1.2)
    group.top = Inches(1.5)
    group.width = Inches(5.6)
    group.height = Inches(2.8)

    info = slide.shapes.add_textbox(Inches(6.9), Inches(2.1), Inches(2.4), Inches(1.4))
    itf = info.text_frame
    itf.text = "grpSp"
    itf.paragraphs[0].runs[0].font.bold = True
    p = itf.add_paragraph()
    p.text = "outer + nested"
    p2 = itf.add_paragraph()
    p2.text = "scale + offset"

    add_notes(slide, "Fixture 05 note text")

    out_name = "fixture-05-group-transform.pptx"
    prs.save(OUT_DIR / out_name)
    return out_name


def make_fixture_06_paragraph_spacing() -> str:
    """Fixture 06: paragraph spacing (spcBef/spcAft/spcPts), buAutoNum, buNone."""
    prs = Presentation()

    # ---------- Slide 1: spcBef / spcAft spacing ----------
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s1, (255, 255, 255))

    tb1 = s1.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.6), Inches(0.65))
    tp = tb1.text_frame.paragraphs[0]
    tp.text = "Fixture 06: Paragraph Spacing"
    tp.runs[0].font.size = Pt(24)
    tp.runs[0].font.bold = True
    tp.runs[0].font.color.rgb = RGBColor(31, 78, 121)

    body1 = s1.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(8.2), Inches(4.5))
    tf1 = body1.text_frame
    tf1.word_wrap = True

    tf1.paragraphs[0].text = "No extra spacing (baseline)"
    tf1.paragraphs[0].runs[0].font.size = Pt(16)

    p2 = tf1.add_paragraph()
    p2.text = "Space-before 18 pt"
    p2.runs[0].font.size = Pt(16)
    p2.space_before = Pt(18)

    p3 = tf1.add_paragraph()
    p3.text = "Space-after 12 pt"
    p3.runs[0].font.size = Pt(16)
    p3.space_after = Pt(12)

    p4 = tf1.add_paragraph()
    p4.text = "Both: spcBef=10 spcAft=10"
    p4.runs[0].font.size = Pt(16)
    p4.space_before = Pt(10)
    p4.space_after = Pt(10)

    add_notes(s1, "Fixture 06 slide 1: paragraph spacing")

    # ---------- Slide 2: buAutoNum auto-numbered list ----------
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s2, (248, 252, 255))

    tb2 = s2.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.6), Inches(0.65))
    tp2 = tb2.text_frame.paragraphs[0]
    tp2.text = "Fixture 06: Auto-numbered List"
    tp2.runs[0].font.size = Pt(24)
    tp2.runs[0].font.bold = True
    tp2.runs[0].font.color.rgb = RGBColor(31, 78, 121)

    body2 = s2.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(8.2), Inches(4.5))
    tf2 = body2.text_frame
    tf2.word_wrap = True

    numbered_items = [
        ("Alpha item", 0),
        ("Beta item", 0),
        ("Gamma item", 0),
        ("Nested sub-item", 1),
        ("Delta item", 0),
    ]
    for idx, (label, lvl) in enumerate(numbered_items):
        para = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
        para.text = label
        para.level = lvl
        para.runs[0].font.size = Pt(16)
        # Inject buAutoNum via XML
        pPr = para._p.get_or_add_pPr()
        ban = etree.SubElement(pPr, qn("a:buAutoNum"))
        ban.set("type", "arabicPeriod")

    add_notes(s2, "Fixture 06 slide 2: auto-numbered list")

    # ---------- Slide 3: buNone + spcPts fixed line height ----------
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(s3, (255, 252, 245))

    tb3 = s3.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(8.6), Inches(0.65))
    tp3 = tb3.text_frame.paragraphs[0]
    tp3.text = "Fixture 06: buNone + Fixed Line Spacing"
    tp3.runs[0].font.size = Pt(24)
    tp3.runs[0].font.bold = True
    tp3.runs[0].font.color.rgb = RGBColor(31, 78, 121)

    body3 = s3.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(8.2), Inches(4.5))
    tf3 = body3.text_frame
    tf3.word_wrap = True

    no_bullet_items = [
        "Fixed line height 22 pt (buNone)",
        "Same spacing, no bullet",
        "Third line, consistent gap",
    ]
    for idx, label in enumerate(no_bullet_items):
        para = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
        para.text = label
        para.runs[0].font.size = Pt(16)
        para.line_spacing = Pt(22)
        # Inject buNone to suppress any inherited bullet
        pPr = para._p.get_or_add_pPr()
        etree.SubElement(pPr, qn("a:buNone"))

    add_notes(s3, "Fixture 06 slide 3: buNone and spcPts line spacing")

    out_name = "fixture-06-paragraph-spacing.pptx"
    prs.save(OUT_DIR / out_name)
    return out_name


def build_manifest(files: list[str]) -> None:
    manifest = {
        "version": 1,
        "generator": "tools/generate_test_pptx.py",
        "fixtures": [
            {
                "file": files[0],
                "purpose": "text + shapes + line + notes",
                "expected": {
                    "slides": 1,
                    "text_min": 2,
                    "shapes_min": 3,
                    "images": 0,
                    "tables": 0,
                    "notes_non_empty": True,
                },
            },
            {
                "file": files[1],
                "purpose": "image + table + notes",
                "expected": {
                    "slides": 1,
                    "text_min": 1,
                    "shapes_min": 0,
                    "images": 1,
                    "tables": 1,
                    "notes_non_empty": True,
                },
            },
            {
                "file": files[2],
                "purpose": "multi-slide + background + notes",
                "expected": {
                    "slides": 3,
                    "text_min": 3,
                    "shapes_min": 0,
                    "images": 0,
                    "tables": 0,
                    "notes_non_empty": True,
                },
            },
            {
                "file": files[3],
                "purpose": "bullets + rotation + notes",
                "expected": {
                    "slides": 1,
                    "text_min": 2,
                    "shapes_min": 1,
                    "images": 0,
                    "tables": 0,
                    "notes_non_empty": True,
                },
            },
            {
                "file": files[4],
                "purpose": "group + nested-group transform + notes",
                "expected": {
                    "slides": 1,
                    "text_min": 3,
                    "shapes_min": 5,
                    "images": 0,
                    "tables": 0,
                    "notes_non_empty": True,
                },
            },
            {
                "file": files[5],
                "purpose": "paragraph spacing + buAutoNum + buNone + spcPts",
                "expected": {
                    "slides": 3,
                    "text_min": 9,
                    "shapes_min": 0,
                    "images": 0,
                    "tables": 0,
                    "notes_non_empty": True,
                },
            },
        ],
    }
    MANIFEST_PATH.write_text(json.dumps(manifest, indent=2), encoding="utf-8")


def main() -> None:
    ensure_dirs()
    image_path = create_sample_image()

    f1 = make_fixture_01_basic_text_shapes()
    f2 = make_fixture_02_image_table(image_path)
    f3 = make_fixture_03_multislide_notes()
    f4 = make_fixture_04_bullets_rotation()
    f5 = make_fixture_05_group_transform()
    f6 = make_fixture_06_paragraph_spacing()

    build_manifest([f1, f2, f3, f4, f5, f6])

    print("Generated fixtures:")
    for p in sorted(OUT_DIR.glob("*.pptx")):
        print("-", p.relative_to(ROOT))
    print("-", MANIFEST_PATH.relative_to(ROOT))


if __name__ == "__main__":
    main()
